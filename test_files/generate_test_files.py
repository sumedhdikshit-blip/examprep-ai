import os
import urllib.request
import docx
from pptx import Presentation

def generate_md():
    print("Generating valid_sample.md...")
    md_content = """# Study Guide: Computer Networks

This is a comprehensive study guide covering the fundamentals of Computer Networks, specifically the OSI Model, TCP/IP Suite, and routing algorithms.

## Section 1: The OSI Model
The Open Systems Interconnection model is a conceptual model that characterizes and standardizes the communication functions of a telecommunication or computing system.

1. Physical Layer: Transmits raw bits over a physical medium.
2. Data Link Layer: Node-to-node data transfer, error detection and correction.
3. Network Layer: Routing of packets across network boundaries (IP addresses).
4. Transport Layer: End-to-end reliability and flow control (TCP, UDP).
5. Session Layer: Managing sessions between applications.
6. Presentation Layer: Translation, encryption, and compression of data.
7. Application Layer: High-level protocols (HTTP, FTP, SMTP).

## Section 2: Routing Algorithms
Routing is the process of selecting paths in a network along which to send network traffic.
- Distance Vector Routing: Uses Bellman-Ford algorithm (e.g., RIP).
- Link State Routing: Uses Dijkstra's algorithm to calculate the shortest path (e.g., OSPF).
"""
    with open("test_files/valid_sample.md", "w", encoding="utf-8") as f:
        f.write(md_content.strip())
    print("valid_sample.md generated successfully.")

def generate_docx():
    print("Generating valid_sample.docx...")
    doc = docx.Document()
    doc.add_heading("Study Guide: Database Management Systems", 0)
    
    # We want to create ~1200 words to test the page approximation
    doc.add_heading("Section 1: Database Normalization", level=1)
    
    p1 = ("Database normalization is the process of structuring a database, usually a relational database, "
          "in accordance with a series of so-called normal forms in order to reduce data redundancy and "
          "improve data integrity. It was first proposed by Edgar F. Codd, an English computer scientist. "
          "Normalization entails organizing the columns and tables of a database to ensure that their dependencies "
          "are properly enforced by database constraints. It is accomplished by applying rules either to synthesis "
          "or decomposition of relationships.")
    doc.add_paragraph(p1)
    
    p2 = ("The primary goal of database normalization is to eliminate database anomalies, which can occur during "
          "insert, update, and delete operations. There are three types of anomalies: Insert Anomalies, where it is "
          "impossible to store certain information unless other information is also stored; Update Anomalies, where "
          "the same data value is stored multiple times in different rows, and updating one row without updating others "
          "creates inconsistency; and Delete Anomalies, where deleting a row accidentally destroys unrelated information "
          "that was only stored in that row.")
    doc.add_paragraph(p2)

    long_text = ("This is a long filler paragraph to increase the word count of the Word document to verify "
                 "that page breaks are approximated correctly. Word count is tracked at paragraph boundaries. "
                 "If a document contains a lot of text, we split it into 500-word logical pages. This allows "
                 "our chunking service to map chunks to approximate page numbers, which is very useful for users. "
                 "Let's write some more content about database keys. A primary key is a unique identifier for a row. "
                 "A foreign key is a column or group of columns that provides a link between data in two tables. "
                 "A candidate key is a column that can uniquely identify a row. A super key is any set of columns "
                 "that uniquely identifies a row. A composite key is a primary key consisting of multiple columns. "
                 "Understanding keys is fundamental to proper database design. In relational database theory, "
                 "a relation is in first normal form if and only if the domain of each attribute contains only atomic "
                 "values, and the value of each attribute contains only a single value from that domain. "
                 "Second normal form requires that the table is in first normal form, and all non-key attributes are "
                 "fully functionally dependent on the primary key, meaning there are no partial dependencies. "
                 "Third normal form requires that the table is in second normal form, and all non-key attributes are "
                 "non-transitively dependent on the primary key. Boyce-Codd normal form is a stronger version of "
                 "third normal form, where for every non-trivial functional dependency X -> Y, X must be a super key. "
                 "Normalizing database schemas up to Boyce-Codd normal form is standard practice for transaction processing.")
    
    for i in range(5):
        doc.add_paragraph(f"Iteration {i+1} of database theory details. " + long_text)
        
    doc.save("test_files/valid_sample.docx")
    print("valid_sample.docx generated successfully.")

def generate_pptx():
    print("Generating valid_sample.pptx...")
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Study Guide: Operating Systems"
    slide1.placeholders[1].text = "Day 2 Pipeline End-to-End Test Presentation"
    
    # Slide 2: Content
    slide_layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide_layout2)
    slide2.shapes.title.text = "Topic 1: Process Scheduling"
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "Process scheduling is the activity of the process manager that handles the removal of the running process from the CPU and the selection of another process."
    p = tf.add_paragraph()
    p.text = "- FIFO (First In First Out): Simplest scheduling algorithm."
    p2 = tf.add_paragraph()
    p2.text = "- SJF (Shortest Job First): Optimal average waiting time."
    p3 = tf.add_paragraph()
    p3.text = "- Round Robin: Designed for time-sharing systems."
    
    # Slide 3: Content
    slide3 = prs.slides.add_slide(slide_layout2)
    slide3.shapes.title.text = "Topic 2: Virtual Memory"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Virtual memory is a memory management capability of an operating system that uses hardware and software to allow a computer to compensate for physical memory shortages."
    p4 = tf3.add_paragraph()
    p4.text = "- Paging: Dividing memory into fixed-size blocks (pages)."
    p5 = tf3.add_paragraph()
    p5.text = "- Page Fault: Occurs when a program accesses a page that is not in physical memory."
    
    prs.save("test_files/valid_sample.pptx")
    print("valid_sample.pptx generated successfully.")

def download_pdf():
    print("Downloading valid_sample.pdf...")
    pdf_url = "https://www.w3.org/WAI/ER/tests/xhtml/testfiles/resources/pdf/dummy.pdf"
    try:
        urllib.request.urlretrieve(pdf_url, "test_files/valid_sample.pdf")
        print("valid_sample.pdf downloaded successfully.")
    except Exception as e:
        print("Failed to download PDF.", e)

def main():
    os.makedirs("test_files", exist_ok=True)
    generate_md()
    generate_docx()
    generate_pptx()
    download_pdf()
    print("All test files generated in 'test_files' directory.")

if __name__ == "__main__":
    main()
