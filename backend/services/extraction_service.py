import os
import pdfplumber
import docx
from pptx import Presentation

def extract_text_from_md(file_path: str) -> list[tuple[int | None, str]]:
    """Read markdown file directly as plain text."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return [(None, content)]
    except Exception as e:
        raise ValueError(f"Failed to read Markdown file: {str(e)}")

def extract_text_from_pdf(file_path: str) -> list[tuple[int | None, str]]:
    """Extract text page by page from PDF using pdfplumber."""
    pages = []
    try:
        with pdfplumber.open(file_path) as pdf:
            if not pdf.pages:
                raise ValueError("PDF file has no pages.")
            for page in pdf.pages:
                text = page.extract_text() or ""
                pages.append((page.page_number, text))
        return pages
    except Exception as e:
        raise ValueError(f"Failed to extract text from PDF: {str(e)}")

def extract_text_from_docx(file_path: str) -> list[tuple[int | None, str]]:
    """Extract paragraph text from DOCX, approximating a page every ~500 words."""
    try:
        doc = docx.Document(file_path)
        pages = []
        current_page_paragraphs = []
        current_word_count = 0
        page_num = 1

        for p in doc.paragraphs:
            p_text = p.text.strip()
            if not p_text:
                continue
            words = p_text.split()
            current_page_paragraphs.append(p_text)
            current_word_count += len(words)

            if current_word_count >= 500:
                pages.append((page_num, "\n".join(current_page_paragraphs)))
                current_page_paragraphs = []
                current_word_count = 0
                page_num += 1

        if current_page_paragraphs:
            pages.append((page_num, "\n".join(current_page_paragraphs)))

        if not pages:
            pages.append((1, ""))
            
        return pages
    except Exception as e:
        raise ValueError(f"Failed to extract text from Word document: {str(e)}")

def extract_text_from_pptx(file_path: str) -> list[tuple[int | None, str]]:
    """Extract text slide by slide from PPTX, using slide index as the page number."""
    try:
        prs = Presentation(file_path)
        pages = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            # slide number is 1-indexed
            pages.append((i + 1, "\n".join(slide_text)))
        
        if not pages:
            pages.append((1, ""))
            
        return pages
    except Exception as e:
        raise ValueError(f"Failed to extract text from PowerPoint presentation: {str(e)}")

def extract_text(file_path: str, file_type: str) -> list[tuple[int | None, str]]:
    """Generic entry point to extract text based on file type."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    # Standardize file_type
    ftype = file_type.lower().strip(".")
    
    if ftype == "md":
        return extract_text_from_md(file_path)
    elif ftype == "pdf":
        return extract_text_from_pdf(file_path)
    elif ftype == "docx":
        return extract_text_from_docx(file_path)
    elif ftype == "pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
